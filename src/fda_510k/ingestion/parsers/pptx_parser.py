from pathlib import Path

from pptx import Presentation

from fda_510k.ingestion.parsers.base import BaseParser, ParsedPage, ParserResult


class PptxParser(BaseParser):
    extensions = (".pptx",)

    def parse(self, path: Path) -> ParserResult:
        prs = Presentation(path)
        pages: list[ParsedPage] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            pages.append(ParsedPage(page_number=i, text="\n".join(texts)))
        return ParserResult(pages=pages)
